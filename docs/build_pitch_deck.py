"""
Generate SyrenBridge Client Pitch Deck (PowerPoint).
Run: python docs/build_pitch_deck.py
Output: docs/SyrenBridge_Pitch_Deck.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
import copy

OUT_PATH = os.path.join(os.path.dirname(__file__), "SyrenBridge_Pitch_Deck.pptx")
LOGO_PATH = os.path.join(os.path.dirname(__file__), "..", "lb_migration_platform_ui", "logo.png")

# ── Colours ───────────────────────────────────────────────────────────────────
ORANGE  = RGBColor(0xFF, 0x36, 0x21)
DARK    = RGBColor(0x09, 0x09, 0x0E)
SLATE   = RGBColor(0x0F, 0x17, 0x2A)
MID     = RGBColor(0x1E, 0x29, 0x3B)
BORDER  = RGBColor(0x33, 0x41, 0x55)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
MUTED   = RGBColor(0x94, 0xA3, 0xB8)
INDIGO  = RGBColor(0x63, 0x66, 0xF1)
CYAN    = RGBColor(0x06, 0xB6, 0xD4)
GREEN   = RGBColor(0x05, 0x96, 0x69)
AMBER   = RGBColor(0xD9, 0x77, 0x06)
LIGHT   = RGBColor(0xF1, 0xF5, 0xF9)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]   # completely blank


# ── Helpers ───────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        if line_w:
            shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             italic=False, wrap=True, font_name="Calibri"):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txb


def add_para(tf, text, size=14, bold=False, color=LIGHT, align=PP_ALIGN.LEFT,
             space_before=0, italic=False):
    from pptx.util import Pt as _Pt
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = _Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def slide_background(slide, color=DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_logo(slide, x=Inches(0.3), y=Inches(0.18), size=Inches(0.55)):
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, x, y, size, size)


def add_footer(slide, text="SyrenBridge by Syren Cloud  |  Confidential  |  2026",
               page_num=None, page_text=None):
    add_rect(slide, 0, H - Inches(0.42), W, Inches(0.42), fill=MID)
    add_text(slide, text,
             Inches(0.4), H - Inches(0.38), Inches(9), Inches(0.32),
             size=9, color=MUTED, align=PP_ALIGN.LEFT)
    if page_text:
        add_text(slide, page_text,
                 Inches(12), H - Inches(0.38), Inches(1.1), Inches(0.32),
                 size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def orange_bar(slide, height=Inches(0.08), y=None):
    y = y if y is not None else Inches(0.9)
    add_rect(slide, 0, y, W, height, fill=ORANGE)


def section_header(slide, title, subtitle=None):
    slide_background(slide, DARK)
    add_rect(slide, 0, 0, W, Inches(1.0), fill=SLATE)
    add_logo(slide)
    add_text(slide, "SYRENBRIDGE", Inches(1.0), Inches(0.22), Inches(4), Inches(0.5),
             size=11, bold=True, color=ORANGE)
    add_text(slide, title, Inches(0.5), Inches(1.2), Inches(12), Inches(1.0),
             size=34, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, Inches(0.5), Inches(2.1), Inches(11), Inches(0.6),
                 size=16, color=MUTED, italic=True)
    orange_bar(slide, height=Inches(0.06), y=Inches(1.0))
    add_footer(slide)


def content_slide_base(slide, title, badge=None, badge_color=ORANGE):
    slide_background(slide, DARK)
    add_rect(slide, 0, 0, W, Inches(1.05), fill=SLATE)
    add_logo(slide)
    add_text(slide, "SYRENBRIDGE", Inches(1.0), Inches(0.22), Inches(4), Inches(0.5),
             size=11, bold=True, color=ORANGE)
    add_text(slide, title, Inches(0.5), Inches(0.18), Inches(10), Inches(0.7),
             size=22, bold=True, color=WHITE)
    if badge:
        add_rect(slide, Inches(11.5), Inches(0.22), Inches(1.5), Inches(0.45), fill=badge_color)
        add_text(slide, badge, Inches(11.5), Inches(0.22), Inches(1.5), Inches(0.45),
                 size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    orange_bar(slide, height=Inches(0.05), y=Inches(1.05))
    add_footer(slide)


def card(slide, l, t, w, h, title, body_lines, title_color=ORANGE,
         bg=MID, border=BORDER, icon=None):
    add_rect(slide, l, t, w, h, fill=bg, line=border, line_w=Pt(0.5))
    cy = t + Inches(0.18)
    if icon:
        add_text(slide, icon, l + Inches(0.15), cy, Inches(0.4), Inches(0.4),
                 size=18, color=title_color)
        add_text(slide, title, l + Inches(0.55), cy, w - Inches(0.7), Inches(0.4),
                 size=13, bold=True, color=title_color)
    else:
        add_text(slide, title, l + Inches(0.2), cy, w - Inches(0.35), Inches(0.4),
                 size=13, bold=True, color=title_color)
    cy += Inches(0.45)
    for line in body_lines:
        add_text(slide, line, l + Inches(0.2), cy, w - Inches(0.35), Inches(0.35),
                 size=10.5, color=LIGHT)
        cy += Inches(0.32)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank_layout)
slide_background(s1, DARK)
add_rect(s1, 0, 0, W, Inches(0.9), fill=ORANGE)
add_text(s1, "SYRENBRIDGE", Inches(0.5), Inches(0.18), Inches(6), Inches(0.55),
         size=15, bold=True, color=WHITE)
add_text(s1, "by Syren Cloud", Inches(0.5), Inches(0.52), Inches(6), Inches(0.35),
         size=11, color=RGBColor(0xFF,0xAA,0x99))

if os.path.exists(LOGO_PATH):
    s1.shapes.add_picture(LOGO_PATH, Inches(11.8), Inches(0.1), Inches(0.7), Inches(0.7))

# Headline
add_text(s1, "Accelerate Your Migration\nto Databricks.", Inches(0.7), Inches(1.6),
         Inches(8), Inches(2.2), size=44, bold=True, color=WHITE)

# Orange accent line
add_rect(s1, Inches(0.7), Inches(3.7), Inches(1.2), Inches(0.07), fill=ORANGE)

add_text(s1, "SyrenBridge automates the conversion of legacy data platforms —\nSQL, ETL, reports, and workflows — to Databricks in days, not months.",
         Inches(0.7), Inches(3.85), Inches(9), Inches(0.9), size=14, color=MUTED)

# Stats row
for i, (val, label, col) in enumerate([
    ("36", "Source Technologies", ORANGE),
    ("13", "Transpiler Dialects", INDIGO),
    ("100%", "Cloud-Native Deployment", CYAN),
]):
    x = Inches(0.7 + i * 3.2)
    add_rect(s1, x, Inches(4.9), Inches(2.9), Inches(1.4), fill=MID, line=BORDER, line_w=Pt(0.5))
    add_text(s1, val, x + Inches(0.15), Inches(4.98), Inches(2.6), Inches(0.65),
             size=36, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s1, label, x + Inches(0.15), Inches(5.55), Inches(2.6), Inches(0.4),
             size=11, color=MUTED, align=PP_ALIGN.CENTER)

add_rect(s1, 0, H - Inches(0.55), W, Inches(0.55), fill=SLATE)
add_text(s1, "Enterprise Migration Platform  ·  Powered by Databricks Labs Lakebridge",
         Inches(0.5), H - Inches(0.48), Inches(10), Inches(0.38),
         size=10, color=MUTED, align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ═════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank_layout)
content_slide_base(s2, "The Migration Challenge", badge="PROBLEM", badge_color=RGBColor(0xDC,0x26,0x26))

problems = [
    ("Legacy ETL is Expensive", ["Maintaining SSIS, Informatica, DataStage, and Oozie", "pipelines costs 3-5x more than modern Spark-native code.", "License costs alone can exceed platform migration costs."]),
    ("Manual Migration Takes Years", ["Typical enterprise DW-to-cloud migration projects", "span 18-36 months — mostly hand-rewriting SQL and ETL.", "Skill shortages in legacy technologies slow every project."]),
    ("Reports & Workflows Left Behind", ["SSRS reports and Oozie workflows are often missed", "in migration scopes, creating shadow BI and", "unsupported orchestration running on old infrastructure."]),
    ("Risk of Getting It Wrong", ["Untested conversions break production pipelines.", "Data quality regressions discovered weeks after cutover.", "No automated way to validate output correctness."]),
]

for i, (title, lines) in enumerate(problems):
    x = Inches(0.4 + (i % 2) * 6.4)
    y = Inches(1.3 + (i // 2) * 2.6)
    add_rect(s2, x, y, Inches(6.0), Inches(2.3), fill=MID, line=BORDER, line_w=Pt(0.5))
    add_rect(s2, x, y, Inches(0.12), Inches(2.3), fill=RGBColor(0xDC,0x26,0x26))
    add_text(s2, title, x + Inches(0.3), y + Inches(0.15), Inches(5.6), Inches(0.45),
             size=14, bold=True, color=WHITE)
    for j, line in enumerate(lines):
        add_text(s2, line, x + Inches(0.3), y + Inches(0.55 + j*0.35), Inches(5.6), Inches(0.35),
                 size=10.5, color=MUTED)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE SOLUTION
# ═════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank_layout)
content_slide_base(s3, "SyrenBridge — The Solution", badge="SOLUTION", badge_color=GREEN)

add_text(s3,
    "SyrenBridge is an enterprise migration platform built on Databricks Apps that automates "
    "assessment, conversion, and validation of legacy data assets — deployed in your Databricks workspace, no additional infrastructure.",
    Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.75), size=13, color=LIGHT)

pills = [
    ("Analyze",    "Complexity scoring for 36 source technologies in minutes",          ORANGE),
    ("Transpile",  "Automated code conversion across 13 dialects to PySpark / SparkSQL", INDIGO),
    ("Validate",   "Built-in SQL validator with dummy data generation",                  GREEN),
    ("Deploy",     "One-click upload of converted artifacts to Databricks workspace",     CYAN),
]
for i, (label, desc, col) in enumerate(pills):
    x = Inches(0.4 + i * 3.15)
    add_rect(s3, x, Inches(2.1), Inches(2.9), Inches(0.5), fill=col)
    add_text(s3, label, x, Inches(2.1), Inches(2.9), Inches(0.5),
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s3, x, Inches(2.6), Inches(2.9), Inches(1.35), fill=MID, line=col, line_w=Pt(0.6))
    add_text(s3, desc, x + Inches(0.12), Inches(2.68), Inches(2.65), Inches(1.2),
             size=10.5, color=LIGHT, wrap=True)

add_text(s3, "How it works:", Inches(0.5), Inches(4.1), Inches(3), Inches(0.45),
         size=14, bold=True, color=ORANGE)
steps = [
    "1  Upload source files or connect to Databricks workspace",
    "2  Run Analyzer — get complexity report and migration readiness score",
    "3  Select dialect in Transpiler — run automated conversion",
    "4  Review output, download ZIP or push directly to workspace",
    "5  Deploy converted notebooks and workflows to production",
]
for j, step in enumerate(steps):
    add_text(s3, step, Inches(0.5), Inches(4.55 + j * 0.42), Inches(12), Inches(0.4),
             size=11.5, color=LIGHT)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — ANALYZER CAPABILITIES
# ═════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank_layout)
content_slide_base(s4, "Analyzer — 36 Source Technologies", badge="ANALYZER", badge_color=CYAN)

add_text(s4, "Instantly assess migration complexity across your entire data estate.",
         Inches(0.5), Inches(1.15), Inches(12), Inches(0.4), size=13, color=MUTED, italic=True)

cats = [
    ("SQL Databases", CYAN, [
        "Athena · BigQuery · Cloudera (Impala) · Greenplum",
        "Hive · IBM DB2 · MS SQL Server · Netezza",
        "Oracle · Presto · Redshift · SAP HANA CalcViews",
        "Snowflake · Synapse · Teradata · Vertica",
    ]),
    ("ETL & Orchestration", AMBER, [
        "ABInitio · ADF · Alteryx · BODS · Datastage",
        "Informatica (PC, BDE, Cloud) · Oozie · PentahoDI",
        "ODI · Sqoop · SSIS · SSRS · Talend",
    ]),
    ("Code & Notebooks", GREEN, [
        "Jupyter Notebooks · PIG · PySpark",
        "SAS · SPSS",
    ]),
]
for i, (title, col, items) in enumerate(cats):
    x = Inches(0.4 + i * 4.3)
    add_rect(s4, x, Inches(1.7), Inches(4.1), Inches(0.48), fill=col)
    add_text(s4, title, x, Inches(1.7), Inches(4.1), Inches(0.48),
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s4, x, Inches(2.18), Inches(4.1), Inches(3.3), fill=MID, line=col, line_w=Pt(0.5))
    for j, item in enumerate(items):
        add_text(s4, item, x + Inches(0.15), Inches(2.28 + j * 0.5), Inches(3.8), Inches(0.45),
                 size=10.5, color=LIGHT)

add_text(s4, "What the Analyzer produces:", Inches(0.5), Inches(5.65), Inches(5), Inches(0.4),
         size=13, bold=True, color=ORANGE)
outputs = [
    "Object inventory — file count, types, line counts",
    "Complexity score per technology — Low / Medium / High",
    "Migration effort estimate (days) per workload",
    "Warnings for unsupported constructs and manual-migration items",
    "Downloadable assessment ZIP with full detail report",
]
for j, o in enumerate(outputs):
    add_text(s4, f"  ✓  {o}", Inches(0.5), Inches(6.05 + j * 0.34), Inches(12), Inches(0.32),
             size=10.5, color=LIGHT)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — TRANSPILER
# ═════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank_layout)
content_slide_base(s5, "Transpiler — 13 Dialects, 3 Engines", badge="TRANSPILER", badge_color=ORANGE)

add_text(s5, "Convert source artifacts to PySpark, SparkSQL, or Databricks Workflow JSON — automatically.",
         Inches(0.5), Inches(1.15), Inches(12), Inches(0.4), size=13, color=MUTED, italic=True)

engines = [
    ("Lakebridge BladeBridge CLI", CYAN, "10 dialects",
     ["DataStage · Informatica · Informatica Cloud",
      "MS SQL Server · Netezza · Oracle",
      "Snowflake · SSIS · Synapse · Teradata"],
     "PySpark or SparkSQL"),
    ("Built-in: HiveSQL Engine", AMBER, "1 dialect",
     ["HiveQL from Cloudera CDH / CDP",
      "sqlglot (read=hive, write=databricks)",
      "Strip Hive-only DDL clauses post-conversion",
      "Optional LLM-assisted fix per statement"],
     "Databricks SQL dialect"),
    ("Built-in: Oozie Engine", AMBER, "1 dialect",
     ["workflow.xml + coordinator.xml",
      "lxml XML parser → Jobs API 2.1 JSON",
      "Coordinator→workflow linking",
      "One-click Databricks Job creation"],
     "Databricks Workflow JSON"),
    ("Built-in: SSRS Engine", AMBER, "1 dialect",
     ["SQL Server Reporting Services .rdl",
      "Extract SQL datasets → .sql notebooks",
      "Flag T-SQL patterns → Spark equivalents",
      "Assessment JSON per report"],
     "SQL Notebooks + Assessment JSON"),
]
for i, (name, col, count, items, output) in enumerate(engines):
    x = Inches(0.35 + i * 3.2)
    add_rect(s5, x, Inches(1.7), Inches(3.0), Inches(0.45), fill=col)
    add_text(s5, name, x + Inches(0.08), Inches(1.73), Inches(2.85), Inches(0.4),
             size=10, bold=True, color=WHITE)
    add_rect(s5, x, Inches(2.15), Inches(3.0), Inches(0.3), fill=SLATE)
    add_text(s5, count, x, Inches(2.15), Inches(3.0), Inches(0.3),
             size=10, color=MUTED, align=PP_ALIGN.CENTER)
    add_rect(s5, x, Inches(2.45), Inches(3.0), Inches(2.5), fill=MID, line=col, line_w=Pt(0.5))
    for j, item in enumerate(items):
        add_text(s5, item, x + Inches(0.12), Inches(2.55 + j * 0.45), Inches(2.8), Inches(0.42),
                 size=9.5, color=LIGHT)
    add_rect(s5, x, Inches(4.95), Inches(3.0), Inches(0.45), fill=DARK, line=col, line_w=Pt(0.5))
    add_text(s5, f"Output: {output}", x + Inches(0.1), Inches(4.97), Inches(2.85), Inches(0.4),
             size=9.5, color=col, bold=True)

add_text(s5,
    "All conversions are performed within your Databricks workspace — source code never leaves your environment.",
    Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.4), size=11.5, color=MUTED, italic=True,
    align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — KEY USE CASES
# ═════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(blank_layout)
content_slide_base(s6, "Key Use Cases", badge="USE CASES", badge_color=INDIGO)

use_cases = [
    ("Data Warehouse Migration", ORANGE,
     "SQL Server / Teradata / Oracle → Databricks SQL",
     ["Convert T-SQL, BTEQ, PL/SQL stored procedures and views",
      "Migrate SSRS reports to Databricks Dashboards (Lakeview)",
      "Migrate SSIS ETL to SparkSQL notebooks"]),
    ("Hadoop Modernisation", CYAN,
     "Cloudera / HDP → Databricks",
     ["Transpile HiveQL to Databricks SQL dialect",
      "Convert Oozie workflows to Databricks Jobs",
      "Migrate Pig Latin scripts with automated mapping"]),
    ("ETL Platform Consolidation", INDIGO,
     "Informatica / DataStage / Talend → Databricks",
     ["Convert PowerCenter mappings to PySpark notebooks",
      "Assess DataStage job complexity before migration",
      "Generate Databricks Workflow JSON from orchestration logic"]),
    ("Cloud-to-Cloud Migration", GREEN,
     "Snowflake / Synapse / Redshift → Databricks",
     ["Transpile modern SQL dialects with high auto-conversion rate",
      "Convert Synapse Pipelines ARM JSON to Databricks Workflows",
      "Preserve CTEs, window functions, and semi-structured types"]),
]

for i, (title, col, subtitle, points) in enumerate(use_cases):
    x = Inches(0.35 + (i % 2) * 6.45)
    y = Inches(1.35 + (i // 2) * 2.85)
    add_rect(s6, x, y, Inches(6.1), Inches(2.65), fill=MID, line=col, line_w=Pt(0.8))
    add_rect(s6, x, y, Inches(6.1), Inches(0.5), fill=col)
    add_text(s6, title, x + Inches(0.15), y + Inches(0.06), Inches(5.8), Inches(0.38),
             size=13, bold=True, color=WHITE)
    add_text(s6, subtitle, x + Inches(0.15), y + Inches(0.52), Inches(5.8), Inches(0.32),
             size=10, color=col, bold=True)
    for j, pt in enumerate(points):
        add_text(s6, f"• {pt}", x + Inches(0.15), y + Inches(0.85 + j * 0.42), Inches(5.8), Inches(0.38),
                 size=10, color=LIGHT)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — HIGH-LEVEL ARCHITECTURE (LAYERED STACK)
# ═════════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(blank_layout)
content_slide_base(s7, "Architecture — How SyrenBridge Works", badge="ARCHITECTURE", badge_color=CYAN)

LAYER_W = Inches(12.4)
LAYER_X = Inches(0.45)

# Layer 1: INPUT
add_rect(s7, LAYER_X, Inches(1.22), LAYER_W, Inches(0.45), fill=SLATE, line=BORDER, line_w=Pt(0.5))
add_text(s7, "INPUT", LAYER_X + Inches(0.15), Inches(1.3), Inches(1.0), Inches(0.3),
         size=9, bold=True, color=MUTED)
for i, (icon, label) in enumerate([
    ("📁", "File Upload / ZIP"),
    ("☁", "Databricks Workspace Browse"),
]):
    bx = LAYER_X + Inches(1.5 + i * 5.4)
    add_rect(s7, bx, Inches(1.26), Inches(4.9), Inches(0.37), fill=MID, line=BORDER, line_w=Pt(0.4))
    add_text(s7, f"{icon}  {label}", bx + Inches(0.15), Inches(1.3), Inches(4.6), Inches(0.28),
             size=10.5, bold=True, color=LIGHT)

add_rect(s7, Inches(6.62), Inches(1.67), Inches(0.08), Inches(0.2), fill=BORDER)

# Layer 2: CORE ENGINES
ENG_Y = Inches(1.87)
ENG_H = Inches(0.9)
add_rect(s7, LAYER_X, ENG_Y, Inches(5.95), ENG_H,
         fill=RGBColor(0x7C, 0x10, 0x0A), line=ORANGE, line_w=Pt(0.6))
add_text(s7, "🔍  ANALYZER", LAYER_X + Inches(0.2), ENG_Y + Inches(0.1), Inches(3.5), Inches(0.36),
         size=13, bold=True, color=WHITE)
add_text(s7, "lakebridge analyze CLI  ·  36 source technologies  ·  Complexity scoring + effort estimate",
         LAYER_X + Inches(0.2), ENG_Y + Inches(0.5), Inches(5.6), Inches(0.32),
         size=9, color=RGBColor(0xFF, 0xAA, 0x99))

add_rect(s7, LAYER_X + Inches(6.45), ENG_Y, Inches(5.95), ENG_H,
         fill=RGBColor(0x2D, 0x2F, 0x8F), line=INDIGO, line_w=Pt(0.6))
add_text(s7, "⚡  TRANSPILER", LAYER_X + Inches(6.65), ENG_Y + Inches(0.1), Inches(3.5), Inches(0.36),
         size=13, bold=True, color=WHITE)
add_text(s7, "Dialect routing  ·  13 dialects  ·  3 engine types  ·  Optional LLM fix per statement",
         LAYER_X + Inches(6.65), ENG_Y + Inches(0.5), Inches(5.6), Inches(0.32),
         size=9, color=RGBColor(0xC7, 0xD2, 0xFE))

add_rect(s7, Inches(9.42), ENG_Y + ENG_H, Inches(0.08), Inches(0.2), fill=BORDER)

# Layer 3: ENGINE DETAIL BOXES
SUB_Y = ENG_Y + ENG_H + Inches(0.2)
SUB_H = Inches(1.05)
sub_w = Inches(3.0)
sub_gap = Inches(0.13)
engines_detail = [
    ("BladeBridge CLI", CYAN,  ["10 dialects", "DataStage · Informatica", "SQL Server · Teradata", "Oracle · Snowflake…"]),
    ("HiveSQL Engine",  AMBER, ["1 dialect", "sqlglot (hive→databricks)", "Clause stripping + rules", "LLM fix (optional)"]),
    ("Oozie Engine",    AMBER, ["1 dialect", "lxml → Jobs API 2.1", "Coordinator linking", "Fan-in DAG preserved"]),
    ("SSRS Engine",     AMBER, ["1 dialect", "RDL → SQL Notebooks", "+ Assessment JSON", "T-SQL pattern flags"]),
]
for i, (name, col, lines) in enumerate(engines_detail):
    sx = LAYER_X + i * (sub_w + sub_gap)
    add_rect(s7, sx, SUB_Y, sub_w, SUB_H, fill=MID, line=col, line_w=Pt(0.5))
    add_rect(s7, sx, SUB_Y, sub_w, Inches(0.3), fill=col)
    add_text(s7, name, sx + Inches(0.1), SUB_Y + Inches(0.05), sub_w - Inches(0.15), Inches(0.24),
             size=9, bold=True, color=WHITE)
    for j, ln in enumerate(lines):
        add_text(s7, ln, sx + Inches(0.12), SUB_Y + Inches(0.36 + j * 0.18), sub_w - Inches(0.2), Inches(0.18),
                 size=8, color=LIGHT)

add_rect(s7, Inches(6.62), SUB_Y + SUB_H, Inches(0.08), Inches(0.2), fill=BORDER)

# Layer 4: OUTPUT
OUT_Y = SUB_Y + SUB_H + Inches(0.2)
add_rect(s7, LAYER_X, OUT_Y, LAYER_W, Inches(0.44),
         fill=RGBColor(0x03, 0x50, 0x38), line=GREEN, line_w=Pt(0.5))
add_text(s7, "OUTPUT", LAYER_X + Inches(0.15), OUT_Y + Inches(0.08), Inches(1.1), Inches(0.3),
         size=9, bold=True, color=MUTED)
add_text(s7,
    "📥 Download ZIP  ·  ☁ Push to Databricks Workspace  ·  "
    "PySpark · SparkSQL · Databricks SQL · Workflow JSON · SQL Notebooks",
    LAYER_X + Inches(1.4), OUT_Y + Inches(0.08), Inches(10.8), Inches(0.3),
    size=9.5, color=RGBColor(0xA7, 0xF3, 0xD0))

add_text(s7,
    "Deployed as Databricks Apps  ·  No external infrastructure  ·  "
    "All conversion runs in-environment  ·  Credentials never leave your workspace",
    LAYER_X, OUT_Y + Inches(0.55), LAYER_W, Inches(0.3),
    size=8.5, color=MUTED, italic=True, align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 7a — DETAILED PROCESS FLOW (USER → WORKSPACE → USER)
# ═════════════════════════════════════════════════════════════════════════════
s7a = prs.slides.add_slide(blank_layout)
content_slide_base(s7a, "End-to-End Process Flow", badge="PROCESS FLOW", badge_color=INDIGO)

add_text(s7a, "Complete user journey — from source file upload to converted artifact running on Databricks.",
         Inches(0.5), Inches(1.15), Inches(12), Inches(0.35), size=11, color=MUTED, italic=True)

CX = Inches(6.63)  # horizontal centre of slide content

def fbox(slide, x, y, w, h, text, sub=None, fill=MID, col=BORDER, tc=WHITE, sc=MUTED, ts=11):
    add_rect(slide, x, y, w, h, fill=fill, line=col, line_w=Pt(0.6))
    ty = y + Inches(0.07) if sub else y + (h - Inches(0.30)) / 2
    add_text(slide, text, x + Inches(0.12), ty, w - Inches(0.22), Inches(0.30),
             size=ts, bold=True, color=tc, align=PP_ALIGN.CENTER)
    if sub:
        add_text(slide, sub, x + Inches(0.12), ty + Inches(0.27), w - Inches(0.22), Inches(0.24),
                 size=8.5, color=sc, align=PP_ALIGN.CENTER)

def vc(slide, cx, y1, y2, col=BORDER):
    if y2 > y1:
        add_rect(slide, cx - Inches(0.035), y1, Inches(0.07), y2 - y1, fill=col)

def hc(slide, x1, x2, cy, col=BORDER):
    if x2 > x1:
        add_rect(slide, x1, cy - Inches(0.035), x2 - x1, Inches(0.07), fill=col)

# ── Row 1: USER START ────────────────────────────────────────────────────────
fbox(s7a, Inches(4.0), Inches(1.20), Inches(5.25), Inches(0.42),
     "👤  User opens SyrenBridge", fill=INDIGO, col=INDIGO, ts=12)
vc(s7a, CX, Inches(1.62), Inches(1.75))
# Branch to two input methods
hc(s7a, Inches(2.4), Inches(10.85), Inches(1.75))
vc(s7a, Inches(2.4),  Inches(1.75), Inches(1.85), col=ORANGE)
vc(s7a, Inches(10.85), Inches(1.75), Inches(1.85), col=CYAN)

# ── Row 2: INPUT METHODS ─────────────────────────────────────────────────────
fbox(s7a, Inches(0.50), Inches(1.85), Inches(4.55), Inches(0.50),
     "📁  Upload Files / ZIP", sub="Local files via browser",
     fill=MID, col=ORANGE, tc=ORANGE)
fbox(s7a, Inches(7.85), Inches(1.85), Inches(4.55), Inches(0.50),
     "☁  Browse Databricks Workspace", sub="Fetch via REST API (Databricks SDK)",
     fill=MID, col=CYAN, tc=CYAN)
# Converge back to centre
vc(s7a, Inches(2.77),  Inches(2.35), Inches(2.48), col=ORANGE)
vc(s7a, Inches(10.12), Inches(2.35), Inches(2.48), col=CYAN)
hc(s7a, Inches(2.77), Inches(10.12), Inches(2.48))
vc(s7a, CX, Inches(2.48), Inches(2.62))

# ── Row 3: APP ROUTING ───────────────────────────────────────────────────────
fbox(s7a, Inches(3.50), Inches(2.62), Inches(6.25), Inches(0.46),
     "⚡  SyrenBridge — Route Request",
     sub="Streamlit App · Deployed on Databricks Apps · No external infra",
     fill=SLATE, col=INDIGO, ts=11)
# Fork to Analyze / Transpile
vc(s7a, CX, Inches(3.08), Inches(3.20))
hc(s7a, Inches(2.85), Inches(10.12), Inches(3.20))
vc(s7a, Inches(2.85),  Inches(3.20), Inches(3.35), col=ORANGE)
vc(s7a, Inches(10.12), Inches(3.20), Inches(3.35), col=INDIGO)

# ── Row 4: BRANCH HEADERS ────────────────────────────────────────────────────
fbox(s7a, Inches(0.50), Inches(3.35), Inches(4.70), Inches(0.36),
     "🔍  ANALYZE", fill=RGBColor(0x7C, 0x10, 0x0A), col=ORANGE, ts=12)
fbox(s7a, Inches(7.40), Inches(3.35), Inches(5.45), Inches(0.36),
     "⚡  TRANSPILE", fill=RGBColor(0x2D, 0x2F, 0x8F), col=INDIGO, ts=12)
vc(s7a, Inches(2.85),  Inches(3.71), Inches(3.82), col=ORANGE)
vc(s7a, Inches(10.12), Inches(3.71), Inches(3.82), col=INDIGO)

# ── Row 5: BRANCH DETAIL ─────────────────────────────────────────────────────
# Analyze detail
add_rect(s7a, Inches(0.50), Inches(3.82), Inches(4.70), Inches(1.20),
         fill=MID, line=ORANGE, line_w=Pt(0.5))
for j, step in enumerate([
    "lakebridge analyze CLI",
    "36 source technologies profiled",
    "Complexity score + effort estimate",
    "Downloadable assessment ZIP",
]):
    add_text(s7a, f"  ▸  {step}", Inches(0.65), Inches(3.93 + j * 0.26), Inches(4.35), Inches(0.24),
             size=9, color=LIGHT)

# 4 engine boxes in transpile branch
EW = Inches(1.22)
EG = Inches(0.12)
EX0 = Inches(7.40)
for i, (ename, ecol, esub) in enumerate([
    ("BladeBridge CLI", CYAN,  "10 dialects"),
    ("HiveSQL Engine",  AMBER, "sqlglot"),
    ("Oozie Engine",    AMBER, "lxml"),
    ("SSRS Engine",     AMBER, "built-in"),
]):
    ex = EX0 + i * (EW + EG)
    add_rect(s7a, ex, Inches(3.82), EW, Inches(1.20), fill=MID, line=ecol, line_w=Pt(0.5))
    add_rect(s7a, ex, Inches(3.82), EW, Inches(0.26), fill=ecol)
    add_text(s7a, ename, ex + Inches(0.05), Inches(3.85), EW - Inches(0.10), Inches(0.22),
             size=7.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s7a, esub, ex + Inches(0.05), Inches(4.14), EW - Inches(0.10), Inches(0.20),
             size=8.0, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(s7a, "→ Artifacts", ex + Inches(0.05), Inches(4.74), EW - Inches(0.10), Inches(0.22),
             size=7.5, color=ecol, align=PP_ALIGN.CENTER)

# ── CONVERGENCE → OUTPUT COLLECTION ─────────────────────────────────────────
MERGE1 = Inches(5.10)
vc(s7a, Inches(2.85),  Inches(5.02), MERGE1, col=ORANGE)
vc(s7a, Inches(10.12), Inches(5.02), MERGE1, col=INDIGO)
hc(s7a, Inches(2.85), Inches(10.12), MERGE1)
vc(s7a, CX, MERGE1, Inches(5.22))
fbox(s7a, Inches(3.20), Inches(5.22), Inches(6.85), Inches(0.46),
     "📊  Results ready — file tree · metrics · preview",
     fill=SLATE, col=GREEN, ts=11)

# Fork to output options
vc(s7a, CX, Inches(5.68), Inches(5.81))
hc(s7a, Inches(3.00), Inches(10.25), Inches(5.81))
vc(s7a, Inches(3.00),  Inches(5.81), Inches(5.92), col=GREEN)
vc(s7a, Inches(10.25), Inches(5.81), Inches(5.92), col=CYAN)

# ── OUTPUT OPTIONS ────────────────────────────────────────────────────────────
fbox(s7a, Inches(0.50), Inches(5.92), Inches(4.90), Inches(0.46),
     "📥  Download ZIP", sub="Converted notebooks / JSONs / SQL files",
     fill=MID, col=GREEN, tc=GREEN)
fbox(s7a, Inches(6.70), Inches(5.92), Inches(5.60), Inches(0.46),
     "☁  Upload to Databricks Workspace", sub="Push artifacts via REST API — no manual copy",
     fill=MID, col=CYAN, tc=CYAN)

# Final convergence
MERGE2 = Inches(6.52)
vc(s7a, Inches(3.00), Inches(6.38), MERGE2, col=GREEN)
vc(s7a, Inches(9.50), Inches(6.38), MERGE2, col=CYAN)
hc(s7a, Inches(3.00), Inches(9.50), MERGE2)
vc(s7a, CX, MERGE2, Inches(6.62))

# ── FINAL RESULT ──────────────────────────────────────────────────────────────
fbox(s7a, Inches(3.50), Inches(6.62), Inches(6.25), Inches(0.38),
     "✅  Converted artifacts ready to run on Databricks",
     fill=RGBColor(0x03, 0x50, 0x38), col=GREEN, ts=11)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 7b — LOW-LEVEL: HiveQL → Databricks SQL EXAMPLE FLOW
# ═════════════════════════════════════════════════════════════════════════════
s7b = prs.slides.add_slide(blank_layout)
content_slide_base(s7b, "Example Flow: HiveQL → Databricks SQL", badge="DEEP DIVE", badge_color=AMBER)

add_text(s7b,
    "Inside the custom HiveSQL engine — converting a Cloudera table definition to Databricks SQL dialect.",
    Inches(0.5), Inches(1.15), Inches(12), Inches(0.38), size=12, color=MUTED, italic=True)

steps_hive = [
    (ORANGE, "1  Upload",             ".hql / .sql file via browser or Databricks workspace browse"),
    (CYAN,   "2  Pre-process",        "Strip comments · detect STORED AS · ROW FORMAT · SERDE · TBLPROPERTIES · LOCATION hdfs://"),
    (CYAN,   "3  Split Statements",   "Line-aware splitter preserves CREATE TABLE … AS SELECT boundaries"),
    (INDIGO, "4  sqlglot",            "transpile(read='hive', write='databricks') applied per statement"),
    (INDIGO, "5  Post-process",       "Regex rules strip surviving Hive DDL · normalise whitespace · fix LIMIT / ORDER"),
    (AMBER,  "6  Issue Tagging",      "Flag constructs needing manual review → attach warning metadata to output"),
    (GREEN,  "7  LLM Fix (optional)", "Per-statement Claude/GPT fix for flagged issues — skipped if LLM not configured"),
]
for j, (col, title, desc) in enumerate(steps_hive):
    y = Inches(1.62 + j * 0.67)
    add_rect(s7b, Inches(0.45), y, Inches(0.07), Inches(0.52), fill=col)
    add_text(s7b, title, Inches(0.66), y + Inches(0.03), Inches(5.3), Inches(0.3),
             size=11, bold=True, color=col)
    add_text(s7b, desc, Inches(0.66), y + Inches(0.3), Inches(5.5), Inches(0.26),
             size=8.5, color=MUTED)
    if j < len(steps_hive) - 1:
        add_rect(s7b, Inches(0.485), y + Inches(0.52), Inches(0.02), Inches(0.15), fill=BORDER)

# Right: BEFORE box
add_rect(s7b, Inches(6.3), Inches(1.55), Inches(6.6), Inches(2.55),
         fill=RGBColor(0x15, 0x10, 0x07), line=AMBER, line_w=Pt(0.5))
add_text(s7b, "BEFORE — HiveQL input (Cloudera CDH)",
         Inches(6.45), Inches(1.62), Inches(6.2), Inches(0.28),
         size=8.5, bold=True, color=AMBER)
before_lines = [
    "CREATE TABLE sales_agg",
    "STORED AS ORC",
    "ROW FORMAT SERDE",
    "  'org.apache.hive.serde2.lazy.LazySimpleSerDe'",
    "TBLPROPERTIES ('orc.compress'='SNAPPY')",
    "LOCATION 'hdfs:///data/sales'",
    "AS SELECT region, SUM(amount) total",
    "FROM raw.transactions WHERE dt='2024-01'",
    "GROUP BY region;",
]
for j, line in enumerate(before_lines):
    add_text(s7b, line, Inches(6.45), Inches(1.93 + j * 0.2), Inches(6.3), Inches(0.2),
             size=7.5, color=RGBColor(0xFB, 0xBF, 0x24), font_name="Courier New")

# Right: AFTER box
add_rect(s7b, Inches(6.3), Inches(4.2), Inches(6.6), Inches(2.0),
         fill=RGBColor(0x03, 0x2A, 0x20), line=GREEN, line_w=Pt(0.5))
add_text(s7b, "AFTER — Databricks SQL output (auto-converted)",
         Inches(6.45), Inches(4.27), Inches(6.2), Inches(0.28),
         size=8.5, bold=True, color=GREEN)
after_lines = [
    "CREATE TABLE sales_agg AS",
    "SELECT region, SUM(amount) total",
    "FROM raw.transactions",
    "WHERE dt = '2024-01'",
    "GROUP BY region;",
]
for j, line in enumerate(after_lines):
    add_text(s7b, line, Inches(6.45), Inches(4.58 + j * 0.23), Inches(6.3), Inches(0.23),
             size=8.5, color=RGBColor(0x6E, 0xE7, 0xB7), font_name="Courier New")

add_text(s7b,
    "✓  STORED AS ORC, ROW FORMAT SERDE, TBLPROPERTIES, LOCATION hdfs:// — all removed automatically",
    Inches(6.45), Inches(5.82), Inches(6.3), Inches(0.28), size=8, color=GREEN)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 7c — SUCCESS METRICS
# ═════════════════════════════════════════════════════════════════════════════
s_metrics = prs.slides.add_slide(blank_layout)
content_slide_base(s_metrics, "Migration Impact — By the Numbers", badge="RESULTS", badge_color=GREEN)

add_text(s_metrics,
    "Outcomes delivered by SyrenBridge-powered Databricks migration programs.",
    Inches(0.5), Inches(1.15), Inches(12), Inches(0.38), size=12, color=MUTED, italic=True)

metrics_data = [
    ("70–90%",        "Reduction in manual\nrewrite effort",                            ORANGE),
    ("Months→Weeks",  "Enterprise migration timeline\ncompressed end-to-end",            CYAN),
    ("36",            "Source technologies assessed\nin minutes, not weeks",             INDIGO),
    ("3–5×",          "Cost reduction vs. maintaining\nlegacy ETL platform licenses",    GREEN),
    ("30 min",        "Deploy SyrenBridge in your\nDatabricks workspace",                AMBER),
    ("13 dialects",   "SQL · ETL · Workflows · Reports\ncovered in one platform",        WHITE),
]
card_w = Inches(3.95)
card_h = Inches(2.3)
for i, (num, label, col) in enumerate(metrics_data):
    mx = Inches(0.4 + (i % 3) * 4.3)
    my = Inches(1.58 + (i // 3) * 2.5)
    add_rect(s_metrics, mx, my, card_w, card_h, fill=MID, line=col, line_w=Pt(0.8))
    add_rect(s_metrics, mx, my, card_w, Inches(0.07), fill=col)
    num_size = 28 if len(num) > 6 else 34
    add_text(s_metrics, num, mx + Inches(0.15), my + Inches(0.14), card_w - Inches(0.25), Inches(0.85),
             size=num_size, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s_metrics, label, mx + Inches(0.15), my + Inches(1.05), card_w - Inches(0.25), Inches(0.85),
             size=10.5, color=LIGHT, align=PP_ALIGN.CENTER, wrap=True)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — COMPETITIVE ADVANTAGE (BATTERY CARDS)
# ═════════════════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(blank_layout)
content_slide_base(s8, "Why SyrenBridge Wins", badge="BATTLE CARDS", badge_color=ORANGE)

add_text(s8,
    "How SyrenBridge compares to every other option in the Databricks migration ecosystem.",
    Inches(0.5), Inches(1.15), Inches(12), Inches(0.38), size=12, color=MUTED, italic=True)

battle_cards = [
    (
        "vs. Lakebridge (Free Open-Source CLI)",
        CYAN,
        [
            "CLI only — no UI, no workspace browser",
            "No Oozie → Databricks Jobs conversion",
            "No HiveSQL Cloudera-specific clause stripping",
            "No SSRS report migration",
        ],
        [
            "Full Streamlit UI — self-serve for any user",
            "Custom Oozie engine: full Jobs API 2.1 JSON",
            "HiveSQL engine: sqlglot + Cloudera clause rules",
            "Built-in SSRS converter (RDL → SQL Notebooks)",
        ],
    ),
    (
        "vs. LeapLogic (Brickbuilder Partner)",
        INDIGO,
        [
            "Oozie: 'assessment only' — no output JSON",
            "Enterprise SI engagement required",
            "Black-box pricing, long onboarding",
            "No SSRS migration documented",
        ],
        [
            "Oozie: full Jobs API 2.1 JSON + coordinator linking",
            "Self-serve — deploy in your workspace in 30 min",
            "Built on open-source Lakebridge — transparent",
            "SSRS → SQL Notebooks included",
        ],
    ),
    (
        "vs. AWS SCT & Cloud-Native Tools",
        AMBER,
        [
            "AWS SCT: Oozie → Step Functions (AWS-only)",
            "No Databricks-native SQL dialect output",
            "Platform lock-in (AWS / Azure / GCP specific)",
            "No SSRS, no Informatica / DataStage migration",
        ],
        [
            "Databricks-first — output runs on SQL Warehouses",
            "Works on AWS, Azure, and GCP Databricks",
            "Oozie → Databricks Jobs (not Step Functions)",
            "Full ETL + SQL + Workflows + Reports in one tool",
        ],
    ),
    (
        "vs. Manual Migration / DIY Teams",
        GREEN,
        [
            "300–400 lines/day manual throughput",
            "18–36 month typical enterprise project timeline",
            "High risk of untested conversion regressions",
            "Requires rare Hive / Oozie specialist skills",
        ],
        [
            "70–90% of rewrite effort automated",
            "Weeks not months — same output, fraction of time",
            "Built-in SQL validator with dummy data generation",
            "No specialist skills needed — works from a browser",
        ],
    ),
]

BCOL_W = Inches(5.95)
BCOL_H = Inches(2.65)
for i, (title, col, them, us) in enumerate(battle_cards):
    bx = Inches(0.35 + (i % 2) * 6.5)
    by = Inches(1.55 + (i // 2) * 2.82)
    add_rect(s8, bx, by, BCOL_W, BCOL_H, fill=MID, line=col, line_w=Pt(0.8))
    add_rect(s8, bx, by, BCOL_W, Inches(0.38), fill=col)
    add_text(s8, title, bx + Inches(0.15), by + Inches(0.06), BCOL_W - Inches(0.25), Inches(0.28),
             size=10.5, bold=True, color=WHITE)
    half_w = (BCOL_W - Inches(0.3)) / 2
    for j, point in enumerate(them):
        add_text(s8, f"✗  {point}",
                 bx + Inches(0.12), by + Inches(0.45 + j * 0.43), half_w, Inches(0.38),
                 size=8, color=RGBColor(0xFC, 0xA5, 0xA5))
    for j, point in enumerate(us):
        add_text(s8, f"✓  {point}",
                 bx + half_w + Inches(0.18), by + Inches(0.45 + j * 0.43), half_w, Inches(0.38),
                 size=8, color=RGBColor(0x86, 0xEF, 0xAC))

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — SSIS + SSRS DEEP DIVE
# ═════════════════════════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(blank_layout)
content_slide_base(s9, "SSIS & SSRS — Deep Dive", badge="SPOTLIGHT", badge_color=AMBER)

# SSIS
add_rect(s9, Inches(0.4), Inches(1.2), Inches(6.1), Inches(5.5), fill=MID, line=CYAN, line_w=Pt(0.8))
add_rect(s9, Inches(0.4), Inches(1.2), Inches(6.1), Inches(0.5), fill=CYAN)
add_text(s9, "SSIS  →  SparkSQL", Inches(0.6), Inches(1.26), Inches(5.8), Inches(0.38),
         size=14, bold=True, color=WHITE)
ssis_points = [
    ("Engine", "Databricks Labs BladeBridge (Lakebridge CLI)"),
    ("Input", ".dtsx, .xml SSIS packages"),
    ("Output", "SparkSQL notebooks (fixed format)"),
    ("Analyzer", "Registered as source #30 — complexity report"),
    ("Converts", "OLE DB Source/Dest, Derived Column, Cond. Split,"),
    ("", "Execute SQL (MERGE), Lookup, Aggregate, Sort"),
    ("Manual", "Script Tasks (C#/VB.NET), Custom components,"),
    ("", "FTP/SMTP tasks, complex @Variable expressions"),
]
for j, (label, val) in enumerate(ssis_points):
    y = Inches(1.82 + j * 0.48)
    if label:
        add_text(s9, label + ":", Inches(0.6), y, Inches(1.4), Inches(0.42), size=10, bold=True, color=CYAN)
    add_text(s9, val, Inches(2.1), y, Inches(4.2), Inches(0.42), size=10, color=LIGHT)

# SSRS
add_rect(s9, Inches(6.8), Inches(1.2), Inches(6.1), Inches(5.5), fill=MID, line=AMBER, line_w=Pt(0.8))
add_rect(s9, Inches(6.8), Inches(1.2), Inches(6.1), Inches(0.5), fill=AMBER)
add_text(s9, "SSRS  →  SQL Notebooks + JSON", Inches(7.0), Inches(1.26), Inches(5.8), Inches(0.38),
         size=14, bold=True, color=WHITE)
ssrs_points = [
    ("Engine", "Built-in ssrs_converter (no CLI required)"),
    ("Input", ".rdl, .rdlc, .rsd reports"),
    ("Output", "One .sql notebook + _assessment.json per report"),
    ("Auto?", "Yes — if report has Text-query datasets"),
    ("Not auto?", "Stored procedure / VB.NET code → JSON only"),
    ("Flags", "GETDATE→current_timestamp, ISNULL→ifnull,"),
    ("", "TOP N→LIMIT N, DATEADD, CONVERT, NOLOCK"),
    ("Manual", "Stored proc migration, VB.NET UDFs, report layout"),
    ("Layout", "Use Databricks Dashboards / Lakeview for visuals"),
]
for j, (label, val) in enumerate(ssrs_points):
    y = Inches(1.82 + j * 0.48)
    if label:
        add_text(s9, label + ":", Inches(7.0), y, Inches(1.5), Inches(0.42), size=10, bold=True, color=AMBER)
    add_text(s9, val, Inches(8.65), y, Inches(4.0), Inches(0.42), size=10, color=LIGHT)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — GETTING STARTED
# ═════════════════════════════════════════════════════════════════════════════
s10 = prs.slides.add_slide(blank_layout)
content_slide_base(s10, "Getting Started", badge="NEXT STEPS", badge_color=GREEN)

add_text(s10, "Deploy SyrenBridge in your Databricks workspace in under 30 minutes.",
         Inches(0.5), Inches(1.15), Inches(12), Inches(0.4), size=13, color=MUTED, italic=True)

steps_detail = [
    ("Step 1", "Deploy to Databricks Apps", GREEN,
     ["Create a Databricks App pointing to the SyrenBridge repo",
      "Install requirements.txt in the app environment",
      "Set DATABRICKS_HOST and DATABRICKS_TOKEN environment variables"]),
    ("Step 2", "Run Your First Analysis", CYAN,
     ["Upload a sample SSIS .dtsx or SQL .sql file",
      "Click 'Run Analyzer' on the Analyzer tab",
      "Review the complexity report and effort estimate"]),
    ("Step 3", "Run Your First Transpilation", ORANGE,
     ["Select your source dialect (e.g., MS SQL Server)",
      "Upload source files and click 'Run Transpiler'",
      "Download the generated PySpark or SparkSQL notebooks"]),
    ("Step 4", "Scale to Production", INDIGO,
     ["Upload full workload ZIPs for batch conversion",
      "Use Databricks workspace browse for direct file access",
      "Upload converted artifacts back to workspace for deployment"]),
]
for i, (step, title, col, points) in enumerate(steps_detail):
    x = Inches(0.35 + (i % 2) * 6.45)
    y = Inches(1.7 + (i // 2) * 2.6)
    add_rect(s10, x, y, Inches(6.1), Inches(2.4), fill=MID, line=col, line_w=Pt(0.6))
    add_rect(s10, x, y, Inches(1.2), Inches(0.5), fill=col)
    add_text(s10, step, x + Inches(0.1), y + Inches(0.07), Inches(1.0), Inches(0.38),
             size=12, bold=True, color=WHITE)
    add_text(s10, title, x + Inches(1.35), y + Inches(0.08), Inches(4.6), Inches(0.38),
             size=12, bold=True, color=col)
    for j, pt in enumerate(points):
        add_text(s10, f"• {pt}", x + Inches(0.2), y + Inches(0.6 + j * 0.42), Inches(5.7), Inches(0.38),
                 size=10, color=LIGHT)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — CLOSING
# ═════════════════════════════════════════════════════════════════════════════
s11 = prs.slides.add_slide(blank_layout)
slide_background(s11, DARK)
add_rect(s11, 0, 0, W, Inches(0.9), fill=ORANGE)
add_text(s11, "SYRENBRIDGE", Inches(0.5), Inches(0.18), Inches(6), Inches(0.55),
         size=15, bold=True, color=WHITE)

if os.path.exists(LOGO_PATH):
    s11.shapes.add_picture(LOGO_PATH, Inches(11.8), Inches(0.1), Inches(0.7), Inches(0.7))

add_text(s11, "Ready to Accelerate Your\nDatabricks Migration?",
         Inches(0.7), Inches(1.6), Inches(9), Inches(2.0), size=40, bold=True, color=WHITE)
add_rect(s11, Inches(0.7), Inches(3.55), Inches(1.2), Inches(0.07), fill=ORANGE)
add_text(s11, "SyrenBridge — 36 source technologies, 13 dialects, deployed in your workspace.",
         Inches(0.7), Inches(3.7), Inches(11), Inches(0.5), size=14, color=MUTED)

# Bottom cards
cards_data = [
    ("Contact", "itsupport@syrencloud.com", ORANGE),
    ("Powered by", "Databricks Labs Lakebridge", CYAN),
    ("Deployed on", "Databricks Apps", INDIGO),
]
for i, (label, val, col) in enumerate(cards_data):
    x = Inches(0.7 + i * 3.9)
    add_rect(s11, x, Inches(4.5), Inches(3.5), Inches(1.2), fill=MID, line=col, line_w=Pt(0.6))
    add_text(s11, label, x + Inches(0.2), Inches(4.6), Inches(3.1), Inches(0.38),
             size=10, color=MUTED)
    add_text(s11, val, x + Inches(0.2), Inches(4.95), Inches(3.1), Inches(0.45),
             size=13, bold=True, color=col)

add_rect(s11, 0, H - Inches(0.55), W, Inches(0.55), fill=SLATE)
add_text(s11, "Syren Cloud  |  2026  |  Confidential",
         Inches(0.5), H - Inches(0.48), Inches(12.3), Inches(0.38),
         size=10, color=MUTED, align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════════
prs.save(OUT_PATH)
print(f"PPT written → {OUT_PATH}")
